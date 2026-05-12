"""
Сборка PPTX-презентации: "Кутузов Михаил Илларионович - великий русский полководец"

Стиль: минимализм + дуотон + крупная типографика + рукописные акценты.
Размер: 16:9 (13.333 x 7.5 inch).
Палитра: фон #F8F8F8, текст #1A1A1A, акцент #0047FF (электрик-синий) и #FF6A00 (неон-оранжевый).

Каждый слайд получает "живой" переход PowerPoint (morph / fade / push / wipe / cover / zoom / split).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree
from copy import deepcopy

# ---------- Палитра ----------
BG      = RGBColor(0xF8, 0xF8, 0xF8)
INK     = RGBColor(0x1A, 0x1A, 0x1A)
MUTE    = RGBColor(0x8A, 0x8A, 0x8A)
LINE    = RGBColor(0xE4, 0xE4, 0xE4)
BLUE    = RGBColor(0x00, 0x47, 0xFF)
NEON    = RGBColor(0xFF, 0x6A, 0x00)
GRAY_GIANT = RGBColor(0xE0, 0xE0, 0xE0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- Шрифты ----------
F_HEAD  = "Montserrat"       # жирный заголовочный
F_BODY  = "Inter"            # основной
F_HAND  = "Segoe Script"     # "рукописный" - есть на большинстве машин
# запасные: PowerPoint подставит подходящий

# ---------- Презентация ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]  # "Blank"


# =========================================================
#           УТИЛИТЫ
# =========================================================

def add_bg(slide, color=BG):
    """Цветной фон слайда."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *,
             size=18, bold=False, color=INK, font=F_BODY,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.1, italic=False):
    """Текстовый блок c единым форматированием."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        f = r.font
        f.name = font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75, radius=None):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shp_type, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    if radius is not None:
        # настраиваем радиус через adjustment (0..0.5)
        try:
            s.adjustments[0] = radius
        except Exception:
            pass
    return s


def add_oval(slide, x, y, w, h, *, fill=None, line=None, line_w=0.75):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    s.shadow.inherit = False
    return s


def add_line(slide, x1, y1, x2, y2, *, color=INK, width=1.5):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_freeform_underline(slide, x, y, w, h, color=NEON, width=6):
    """Волнистое 'маркерное' подчёркивание."""
    ff = slide.shapes.build_freeform(x, y + h/2, scale=1.0)
    # ломаная, имитирующая живой штрих
    steps = 6
    for i in range(1, steps + 1):
        t = i / steps
        px = x + w * t
        # синусоида
        import math
        py = y + h/2 + math.sin(t * 3.14 * 2) * h * 0.45
        ff.add_line_segments([(px, py)], close=False)
    shape = ff.convert_to_shape()
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(width)
    # закруглённые концы
    ln = shape.line._get_or_add_ln()
    ln.set("cap", "rnd")
    return shape


# ----- Переходы между слайдами (XML-инъекция) -----
TRANSITIONS = {
    "fade":   '<p:fade xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>',
    "morph":  '<p:morph xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" option="byObject"/>',
    "push":   '<p:push xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" dir="l"/>',
    "wipe":   '<p:wipe xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" dir="l"/>',
    "cover":  '<p:cover xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" dir="d"/>',
    "split":  '<p:split xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" orient="horz" dir="out"/>',
    "zoom":   '<p:zoom xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>',
    "random": '<p:randomBar xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" dir="vert"/>',
}

def set_transition(slide, kind="fade", duration_ms=900, advance_click=True):
    """Добавляет тег <p:transition> в XML слайда. Включает современный переход (p14:*) где это уместно."""
    nsmap = {
        "p":   "http://schemas.openxmlformats.org/presentationml/2006/main",
        "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
        "mc":  "http://schemas.openxmlformats.org/markup-compatibility/2006",
    }
    sld = slide._element  # <p:sld>
    # удаляем старый transition, если есть
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)

    trans = etree.SubElement(
        sld,
        qn("p:transition"),
        attrib={"spd": "med", "advClick": "1" if advance_click else "0"},
        nsmap=None,
    )
    # явно укажем длительность в миллисекундах (p14)
    trans.set("{%s}dur" % nsmap["p14"], str(duration_ms))

    effect_xml = TRANSITIONS.get(kind, TRANSITIONS["fade"])
    effect = etree.fromstring(effect_xml)
    trans.append(effect)
    return trans


# ----- Нумерация слайда / тег -----
def slide_header(slide, number, total, tag_text):
    # цветная точка-булет
    add_oval(slide, Inches(0.55), Inches(0.42), Inches(0.12), Inches(0.12), fill=BLUE)
    add_text(slide, Inches(0.72), Inches(0.34), Inches(6), Inches(0.3),
             tag_text.upper(), size=10, color=MUTE, font=F_BODY, bold=True)
    add_text(slide, Inches(12.2), Inches(0.34), Inches(1), Inches(0.3),
             f"{number:02d} / {total:02d}", size=10, color=MUTE, font=F_HEAD,
             align=PP_ALIGN.RIGHT, bold=True)


# =========================================================
#                  СЛАЙДЫ
# =========================================================

TOTAL = 10


# ---------- 1. Титульный ----------
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 1, TOTAL, "Slide 01  -  Title")

    # Левая часть (50%)
    add_text(s, Inches(0.65), Inches(1.1), Inches(6.5), Inches(0.5),
             "Он проиграл Москву,  чтобы выиграть войну.",
             size=20, color=BLUE, font=F_HAND, italic=True)

    # Главный заголовок
    add_text(s, Inches(0.6), Inches(1.7), Inches(6.8), Inches(3.1),
             "Кутузов.\nМихаил\nИлларионович.",
             size=64, bold=True, color=INK, font=F_HEAD, line_spacing=0.95)

    # Маркерное подчёркивание под ключевым словом
    add_rect(s, Inches(0.7), Inches(4.55), Inches(4.2), Inches(0.22),
             fill=NEON)  # плоский маркер
    # Ломаный штрих поверх (имитация "от руки")
    add_freeform_underline(s, Inches(0.7), Inches(4.6), Inches(4.2), Inches(0.18),
                           color=NEON, width=5)

    # Подзаголовок
    add_text(s, Inches(0.65), Inches(5.05), Inches(7), Inches(1.0),
             "Великий русский полководец - или главный мастер\nстратегического отступления в истории Европы?",
             size=16, color=INK, font=F_BODY, line_spacing=1.35)

    # Мета
    add_text(s, Inches(0.65), Inches(6.6), Inches(7), Inches(0.4),
             "1745  -  1813      ·      FIELDMARSHAL      ·      BATTLE OF BORODINO",
             size=10, color=MUTE, font=F_HEAD, bold=True)

    # Правая часть (50%): дуотонный "портрет-абстракция"
    RX = Inches(7.0); RY = Inches(0); RW = Inches(6.333); RH = Inches(7.5)
    add_rect(s, RX, RY, RW, RH, fill=INK)

    # диагональные полосы (паттерн)
    import math
    for i in range(-6, 30):
        x = RX + Inches(i * 0.35)
        add_line(s, x, RY, x + Inches(2), RY + Inches(7.5),
                 color=RGBColor(0x2A, 0x2A, 0x2A), width=0.5)

    # треуголка
    tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                             Inches(8.0), Inches(1.2), Inches(4.3), Inches(1.2))
    tri.rotation = 0
    tri.fill.solid(); tri.fill.fore_color.rgb = BLUE
    tri.line.fill.background(); tri.shadow.inherit = False

    # шляпа - нижняя дуга (трапеция)
    hat = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID,
                             Inches(7.9), Inches(1.9), Inches(4.5), Inches(0.7))
    hat.rotation = 180
    hat.fill.solid(); hat.fill.fore_color.rgb = BLUE
    hat.line.fill.background(); hat.shadow.inherit = False

    # перо
    feather = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                                 Inches(11.5), Inches(0.7), Inches(1.1), Inches(1.2))
    feather.rotation = 20
    feather.fill.solid(); feather.fill.fore_color.rgb = NEON
    feather.line.fill.background(); feather.shadow.inherit = False

    # голова
    add_oval(s, Inches(8.7), Inches(2.4), Inches(3.0), Inches(3.4), fill=BLUE)

    # повязка на глазу
    add_rect(s, Inches(8.5), Inches(3.5), Inches(3.4), Inches(0.5), fill=INK)
    add_line(s, Inches(8.5), Inches(3.55), Inches(11.9), Inches(3.55),
             color=NEON, width=1.5)
    add_line(s, Inches(8.5), Inches(3.95), Inches(11.9), Inches(3.95),
             color=NEON, width=1.5)

    # мундир (трапеция вниз)
    coat = s.shapes.add_shape(MSO_SHAPE.TRAPEZOID,
                              Inches(7.5), Inches(5.6), Inches(5.5), Inches(1.9))
    coat.fill.solid(); coat.fill.fore_color.rgb = BLUE
    coat.line.fill.background(); coat.shadow.inherit = False

    # эполеты
    add_oval(s, Inches(7.7), Inches(5.5), Inches(0.9), Inches(0.9), fill=NEON)
    add_oval(s, Inches(11.9), Inches(5.5), Inches(0.9), Inches(0.9), fill=NEON)

    # орден (крест)
    add_rect(s, Inches(10.05), Inches(6.25), Inches(0.2), Inches(0.9), fill=BG)
    add_rect(s, Inches(9.7),  Inches(6.6),  Inches(0.9), Inches(0.2), fill=BG)

    # рукописная заметка
    add_text(s, Inches(7.2), Inches(6.75), Inches(5.2), Inches(0.6),
             "«не победит тот, кто не умеет вовремя уйти» - заметка на полях",
             size=14, color=WHITE, font=F_HAND, italic=True)

    # стрелочка "тот самый"
    add_text(s, Inches(11.6), Inches(1.1), Inches(1.5), Inches(0.4),
             "тот самый  ↙", size=13, color=NEON, font=F_HAND)

    set_transition(s, "fade", duration_ms=800)


# ---------- 2. Проблема ----------
def slide_problem():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 2, TOTAL, "Slide 02  -  Проблема")

    # Гигантская цифра
    add_text(s, Inches(0), Inches(0.6), Inches(13.333), Inches(6),
             "610 000",
             size=400, bold=True, color=GRAY_GIANT, font=F_HEAD,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.85)

    # Поверх - ключевая фраза
    add_text(s, Inches(0.8), Inches(3.0), Inches(7.5), Inches(1.2),
             "Армия Наполеона\nидёт на Россию.",
             size=44, bold=True, color=INK, font=F_HEAD, line_spacing=1.0)

    add_text(s, Inches(0.8), Inches(4.8), Inches(7.0), Inches(1.2),
             "Крупнейшее вторжение в истории Европы. Прямой\nбой - самоубийство. Отступать - позор.",
             size=14, color=INK, font=F_BODY, line_spacing=1.4)

    # штамп от руки
    stamp = add_rect(s, Inches(9.0), Inches(5.8), Inches(3.9), Inches(0.8),
                     fill=None, line=NEON, line_w=1.5, radius=0.3)
    # делаем пунктирную
    ln = stamp.line._get_or_add_ln()
    for child in list(ln):
        if child.tag == qn("a:prstDash"):
            ln.remove(child)
    prstDash = etree.SubElement(ln, qn("a:prstDash"))
    prstDash.set("val", "dash")

    add_text(s, Inches(9.0), Inches(5.8), Inches(3.9), Inches(0.8),
             "Что делать, Михаил Илларионович?",
             size=16, color=NEON, font=F_HAND,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, italic=True)

    set_transition(s, "push", duration_ms=900)


# ---------- 3–5. Слайды-переходы ----------
def slide_break(num, variant, icon_draw):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, num, TOTAL, f"Slide {num:02d}  -  Пауза")

    # цветной блоб сзади
    if variant == "a":
        add_oval(s, Inches(-1.2), Inches(0.5), Inches(4.5), Inches(4.5), fill=BLUE)
        add_oval(s, Inches(10.5), Inches(4.2), Inches(3.5), Inches(3.5), fill=NEON)
    elif variant == "b":
        add_oval(s, Inches(9.5), Inches(-1), Inches(4.5), Inches(4.5), fill=NEON)
        add_oval(s, Inches(-1), Inches(4), Inches(3.5), Inches(3.5), fill=BLUE)
    else:
        add_oval(s, Inches(4.5), Inches(-2), Inches(4.5), Inches(4.5), fill=BLUE)

    # размытие через полупрозрачный слой BG
    overlay = add_rect(s, 0, 0, SW, SH, fill=BG)
    # полупрозрачность через XML
    sp = overlay.fill.fore_color._xFill
    solid = overlay.fill._xPr.find(qn("a:solidFill"))
    if solid is not None:
        srgb = solid.find(qn("a:srgbClr"))
        alpha = etree.SubElement(srgb, qn("a:alpha"))
        alpha.set("val", "65000")  # 65%

    # 3D-сфера (концентрические круги для "объёма")
    cx, cy = Inches(6.666), Inches(3.75)
    r = Inches(2.2)
    add_oval(s, cx - r, cy - r, 2*r, 2*r, fill=RGBColor(0xBF, 0xBF, 0xBF))
    add_oval(s, cx - r + Inches(0.15), cy - r + Inches(0.1),
             2*r - Inches(0.3), 2*r - Inches(0.3),
             fill=RGBColor(0xD6, 0xD6, 0xD6))
    add_oval(s, cx - r + Inches(0.35), cy - r + Inches(0.25),
             2*r - Inches(0.7), 2*r - Inches(0.7),
             fill=RGBColor(0xEE, 0xEE, 0xEE))
    # блик
    add_oval(s, cx - Inches(0.9), cy - Inches(1.4),
             Inches(0.9), Inches(0.55),
             fill=WHITE)

    # иконка в центре
    icon_box_w = Inches(1.4); icon_box_h = Inches(1.4)
    add_rect(s, cx - icon_box_w/2, cy - icon_box_h/2, icon_box_w, icon_box_h,
             fill=INK, radius=0.2)
    add_rect(s, cx - icon_box_w/2, cy - icon_box_h/2 + Inches(1.3),
             icon_box_w, Inches(0.1),
             fill=(BLUE if variant != "b" else NEON), radius=0.3)

    icon_draw(s, cx, cy)

    set_transition(s, {"a":"wipe","b":"cover","c":"split"}[variant], duration_ms=700)


def icon_map(s, cx, cy):
    # грубая карта: трапеция
    add_rect(s, cx - Inches(0.45), cy - Inches(0.35),
             Inches(0.9), Inches(0.7), fill=None, line=WHITE, line_w=2)
    add_line(s, cx - Inches(0.15), cy - Inches(0.35),
             cx - Inches(0.15), cy + Inches(0.35), color=WHITE, width=2)
    add_line(s, cx + Inches(0.15), cy - Inches(0.35),
             cx + Inches(0.15), cy + Inches(0.35), color=WHITE, width=2)


def icon_clock(s, cx, cy):
    add_oval(s, cx - Inches(0.4), cy - Inches(0.4),
             Inches(0.8), Inches(0.8),
             fill=None, line=WHITE, line_w=2)
    add_line(s, cx, cy, cx, cy - Inches(0.3), color=WHITE, width=2)
    add_line(s, cx, cy, cx + Inches(0.22), cy, color=WHITE, width=2)


def icon_snow(s, cx, cy):
    for dx, dy in [(0, 0.4), (0, -0.4), (0.4, 0), (-0.4, 0),
                   (0.28, 0.28), (-0.28, -0.28), (0.28, -0.28), (-0.28, 0.28)]:
        add_line(s, cx, cy, cx + Inches(dx), cy + Inches(dy), color=WHITE, width=2)


# ---------- 6. Решение - бинго ----------
def slide_solution():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 6, TOTAL, "Slide 06  -  Решение")

    add_text(s, Inches(0.65), Inches(1.0), Inches(12), Inches(1.3),
             "Четыре хода Кутузова,\nкоторые заменили «лобовой удар».",
             size=34, bold=True, color=INK, font=F_HEAD, line_spacing=1.05)

    add_text(s, Inches(0.65), Inches(2.65), Inches(11), Inches(0.7),
             "Вместо того чтобы играть по правилам Наполеона, он поменял саму игру:\nвремя, пространство, снабжение и мораль.",
             size=13, color=INK, font=F_BODY, line_spacing=1.35)

    # 4 ячейки
    cells = [
        ("Бородино, 1812", "~ «Генеральное сражение ради сражения»",
         "→ Изматывание ценой поля, не армии", BLUE, "ход №1"),
        ("Тарутинский манёвр", "~ «Оборонять Москву любой ценой»",
         "→ Отдать город, перекрыть юг и снабжение", NEON, "ход №2"),
        ("Пустая Москва", "~ «Склады и зимние квартиры»",
         "→ Пепел вместо ресурсов для врага", NEON, "ход №3"),
        ("Малоярославец + партизаны", "~ «Парадное преследование»",
         "→ Гнать по выжженной Смоленской дороге", BLUE, "ход №4"),
    ]

    GX = Inches(0.65); GY = Inches(3.8)
    CW = Inches(6.0);  CH = Inches(1.65)
    GAP_X = Inches(0.15); GAP_Y = Inches(0.2)

    for i, (h, old, new, accent, corner) in enumerate(cells):
        col = i % 2; row = i // 2
        x = GX + col * (CW + GAP_X)
        y = GY + row * (CH + GAP_Y)
        add_rect(s, x, y, CW, CH, fill=WHITE, line=LINE, line_w=0.5, radius=0.06)
        # цветная полоса сверху
        add_rect(s, x, y, CW, Inches(0.08), fill=accent, radius=0)
        # иконка
        add_rect(s, x + Inches(0.3), y + Inches(0.25),
                 Inches(0.5), Inches(0.5), fill=INK, radius=0.25)
        # заголовок
        add_text(s, x + Inches(1.0), y + Inches(0.2),
                 CW - Inches(1.3), Inches(0.45),
                 h, size=16, bold=True, color=INK, font=F_HEAD)
        # старое (зачёркнуто)
        old_tb = add_text(s, x + Inches(1.0), y + Inches(0.72),
                          CW - Inches(1.3), Inches(0.35),
                          old, size=10, color=MUTE, font=F_BODY)
        # зачёркивание через XML
        for p in old_tb.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.get_or_add_rPr()
                rPr.set("strike", "sngStrike")
        # новое
        add_text(s, x + Inches(1.0), y + Inches(1.05),
                 CW - Inches(1.3), Inches(0.5),
                 new, size=11, bold=True, color=INK, font=F_BODY)
        # рукописный уголок
        add_text(s, x + CW - Inches(1.3), y + Inches(0.15),
                 Inches(1.2), Inches(0.3),
                 corner, size=12, color=accent, font=F_HAND,
                 align=PP_ALIGN.RIGHT, italic=True)

    set_transition(s, "zoom", duration_ms=900)


# ---------- 7. Кейс - имитация чата ----------
def slide_case():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 7, TOTAL, "Slide 07  -  Кейс")

    # Левая колонка - текст
    add_text(s, Inches(0.65), Inches(1.1), Inches(6), Inches(0.4),
             "1812, личная переписка*", size=16, color=NEON, font=F_HAND, italic=True)
    add_text(s, Inches(0.65), Inches(1.6), Inches(6.3), Inches(2.5),
             "Результат,\nкоторого никто не ждал.",
             size=38, bold=True, color=INK, font=F_HEAD, line_spacing=1.0)
    add_text(s, Inches(0.65), Inches(4.4), Inches(6.2), Inches(1.5),
             "* реконструкция по мотивам донесений.\nАрмия Наполеона ушла из России с ~5% от\nпервоначальной численности. Без единого\n«решающего сражения».",
             size=13, color=INK, font=F_BODY, line_spacing=1.4)

    # чипы
    chips = ["#Бородино", "#Тарутино", "#Березина", "-580 000 солдат"]
    cx = Inches(0.65); cy = Inches(6.3)
    for c in chips:
        cw = Inches(0.18 + 0.14 * len(c))
        add_rect(s, cx, cy, cw, Inches(0.4), fill=WHITE, line=LINE, line_w=0.5, radius=0.4)
        add_text(s, cx, cy, cw, Inches(0.4), c, size=10, color=INK, font=F_BODY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + Inches(0.1)

    # Правая колонка - чат
    CHX = Inches(7.3); CHY = Inches(1.2); CHW = Inches(5.5); CHH = Inches(5.6)
    add_rect(s, CHX, CHY, CHW, CHH, fill=WHITE, line=LINE, line_w=0.5, radius=0.05)
    add_text(s, CHX + Inches(0.3), CHY - Inches(0.4),
             CHW, Inches(0.3),
             "TELEGRAM  ·  #СТАВКА 1812",
             size=10, color=MUTE, font=F_HEAD, bold=True)

    def msg(y, avatar_color, avatar_letter, name, name_color, bubble_text,
            bubble_color, time_text, from_me=False, reaction=None):
        if from_me:
            av_x = CHX + CHW - Inches(0.75)
            name_x = CHX + CHW - Inches(4.4)
            name_align = PP_ALIGN.RIGHT
            bubble_x = CHX + CHW - Inches(4.4)
            bubble_w = Inches(3.6)
        else:
            av_x = CHX + Inches(0.3)
            name_x = CHX + Inches(1.05)
            name_align = PP_ALIGN.LEFT
            bubble_x = CHX + Inches(1.05)
            bubble_w = Inches(3.6)

        add_oval(s, av_x, y, Inches(0.55), Inches(0.55), fill=avatar_color)
        add_text(s, av_x, y, Inches(0.55), Inches(0.55),
                 avatar_letter, size=14, bold=True, color=WHITE, font=F_HEAD,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_text(s, name_x, y, Inches(3.6), Inches(0.3),
                 name, size=10, bold=True, color=name_color, font=F_BODY,
                 align=name_align)

        # bubble
        add_rect(s, bubble_x, y + Inches(0.3), bubble_w, Inches(0.9),
                 fill=bubble_color, radius=0.2)
        add_text(s, bubble_x + Inches(0.15), y + Inches(0.38),
                 bubble_w - Inches(0.3), Inches(0.8),
                 bubble_text, size=11, color=INK, font=F_BODY, line_spacing=1.3)

        # time
        add_text(s, bubble_x, y + Inches(1.22), bubble_w, Inches(0.25),
                 time_text + "   ✓✓", size=9, color=MUTE, font=F_BODY,
                 align=name_align)

        if reaction:
            add_rect(s, bubble_x, y + Inches(1.45), Inches(3.2), Inches(0.28),
                     fill=WHITE, line=LINE, line_w=0.5, radius=0.5)
            add_text(s, bubble_x + Inches(0.15), y + Inches(1.47),
                     Inches(3.0), Inches(0.25),
                     reaction, size=9, color=INK, font=F_BODY, align=PP_ALIGN.LEFT)

    msg(CHY + Inches(0.2), BLUE, "А1", "Александр I", BLUE,
        "Михаил Илларионович, вы правда\nотдали Москву?! Это как вообще.",
        RGBColor(0xF2, 0xF4, 0xF7), "14:02", from_me=False)

    msg(CHY + Inches(2.1), NEON, "К", "Кутузов", NEON,
        "Потерей Москвы ещё не потеряна\nРоссия. С потерей армии - Россия\nпогибла.",
        RGBColor(0xE6, 0xEE, 0xFF), "19:47", from_me=True,
        reaction="🔥 Барклай   🫡 Ермолов   ❤ Давыдов")

    msg(CHY + Inches(4.4), INK, "N", "N. Bonaparte", INK,
        "прочитано  ✓✓",
        RGBColor(0xFF, 0xEC, 0xE0), "—", from_me=False)

    set_transition(s, "cover", duration_ms=900)


# ---------- 8. Финал с QR ----------
def slide_final():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 10, TOTAL, "Slide 10  -  Финал")

    # Заголовок
    add_text(s, Inches(0.65), Inches(1.6), Inches(8), Inches(0.6),
             "давайте поговорим →", size=18, color=MUTE, font=F_HAND, italic=True)

    # "Давайте сделаем это" - "сделаем" оранжевым
    add_text(s, Inches(0.65), Inches(2.1), Inches(8.5), Inches(1.8),
             "Давайте", size=80, bold=True, color=INK, font=F_HEAD, line_spacing=0.95)
    add_text(s, Inches(0.65), Inches(3.2), Inches(8.5), Inches(1.8),
             "сделаем", size=80, bold=True, color=NEON, font=F_HEAD, line_spacing=0.95)
    add_text(s, Inches(0.65), Inches(4.3), Inches(8.5), Inches(1.8),
             "это.", size=80, bold=True, color=INK, font=F_HEAD, line_spacing=0.95)

    add_text(s, Inches(0.65), Inches(5.7), Inches(7.5), Inches(1.0),
             "Вопрос не в том, дадим ли мы бой.\nВопрос - на чьём поле, в чью зиму и по чьим правилам.",
             size=13, color=INK, font=F_BODY, line_spacing=1.4)

    # CTA
    add_rect(s, Inches(0.65), Inches(6.7), Inches(2.8), Inches(0.55),
             fill=NEON, radius=0.2)
    add_oval(s, Inches(0.85), Inches(6.88), Inches(0.18), Inches(0.18), fill=WHITE)
    add_text(s, Inches(1.1), Inches(6.7), Inches(2.6), Inches(0.55),
             "Открыть обсуждение  →",
             size=13, bold=True, color=WHITE, font=F_BODY,
             anchor=MSO_ANCHOR.MIDDLE)

    # QR в пиксель-арт стиле
    QX = Inches(9.0); QY = Inches(1.6); QSIZE = Inches(3.8)
    add_rect(s, QX - Inches(0.25), QY - Inches(0.25),
             QSIZE + Inches(0.5), QSIZE + Inches(0.5),
             fill=WHITE, line=LINE, line_w=0.5, radius=0.05)

    # Подпись над QR
    add_text(s, QX - Inches(0.2), QY - Inches(0.6),
             QSIZE, Inches(0.4),
             "наведи камеру, солдат ✎",
             size=14, color=BLUE, font=F_HAND, italic=True, align=PP_ALIGN.CENTER)

    # Матрица 25x25 - декоративная, с финдерами
    N = 25
    cell = QSIZE / N
    # Базовый паттерн (псевдослучайный, но детерминированный)
    import hashlib
    def fill_cell(r, c):
        h = hashlib.md5(f"{r},{c}".encode()).digest()
        return h[0] % 2 == 0

    for r in range(N):
        for c in range(N):
            # пропускаем области финдеров - заполним отдельно
            in_finder = (
                (r < 7 and c < 7) or
                (r < 7 and c >= N - 7) or
                (r >= N - 7 and c < 7)
            )
            if in_finder:
                continue
            if fill_cell(r, c):
                x = QX + c * cell
                y = QY + r * cell
                add_rect(s, x, y, cell, cell, fill=INK)

    # 3 финдера
    def finder(fx, fy):
        add_rect(s, fx, fy, 7 * cell, 7 * cell, fill=INK)
        add_rect(s, fx + cell, fy + cell, 5 * cell, 5 * cell, fill=WHITE)
        add_rect(s, fx + 2 * cell, fy + 2 * cell, 3 * cell, 3 * cell, fill=INK)

    finder(QX, QY)
    finder(QX + (N - 7) * cell, QY)
    finder(QX, QY + (N - 7) * cell)

    # акцент-логотип в центре
    add_rect(s, QX + 10 * cell, QY + 10 * cell, 5 * cell, 5 * cell, fill=BG)
    add_rect(s, QX + 11 * cell, QY + 11 * cell, 3 * cell, 3 * cell, fill=NEON)

    # Подпись
    add_text(s, Inches(0.65), Inches(7.1), Inches(8), Inches(0.3),
             "ПРЕЗЕНТАЦИЯ  ·  КУТУЗОВ М.И.  ·  1745 - 1813",
             size=9, color=MUTE, font=F_HEAD, bold=True)

    set_transition(s, "morph", duration_ms=1100)


# ---------- 9. Наследие (бонус для плавности повествования) ----------
def slide_legacy():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, BG)
    slide_header(s, 9, TOTAL, "Slide 09  -  Наследие")

    add_text(s, Inches(0.65), Inches(1.1), Inches(12), Inches(0.5),
             "через 210 лет", size=18, color=BLUE, font=F_HAND, italic=True)

    add_text(s, Inches(0.65), Inches(1.7), Inches(12), Inches(2.0),
             "Его имя носят улицы,\nпроспекты и военные доктрины.",
             size=40, bold=True, color=INK, font=F_HEAD, line_spacing=1.05)

    # 3 цифры
    stats = [
        ("1812", "год, после которого\nРоссия стала другой"),
        ("16",   "кампаний в биографии,\nот Крыма до Европы"),
        ("∞",    "страниц учебников\nи военных академий"),
    ]
    bx = Inches(0.65); by = Inches(4.6); bw = Inches(4.0); bh = Inches(2.2)
    for i, (big, sub) in enumerate(stats):
        x = bx + i * (bw + Inches(0.15))
        add_rect(s, x, by, bw, bh, fill=WHITE, line=LINE, line_w=0.5, radius=0.04)
        # цветная полоска
        accent = [BLUE, NEON, BLUE][i]
        add_rect(s, x, by, Inches(0.1), bh, fill=accent)
        add_text(s, x + Inches(0.35), by + Inches(0.2),
                 bw - Inches(0.5), Inches(1.2),
                 big, size=60, bold=True, color=INK, font=F_HEAD, line_spacing=0.9)
        add_text(s, x + Inches(0.35), by + Inches(1.45),
                 bw - Inches(0.5), Inches(0.7),
                 sub, size=12, color=INK, font=F_BODY, line_spacing=1.35)

    set_transition(s, "push", duration_ms=900)


# ---------- 8. Цитата (бонус - переходный "глоток воздуха") ----------
def slide_quote():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, INK)
    slide_header(s, 8, TOTAL, "Slide 08  -  Цитата")

    # большая кавычка
    add_text(s, Inches(0.65), Inches(0.9), Inches(3), Inches(3),
             "“", size=300, bold=True, color=NEON, font=F_HEAD, line_spacing=0.8)

    add_text(s, Inches(2.5), Inches(2.4), Inches(10.5), Inches(3.5),
             "Со взятием Москвы\nещё не покорена Россия.",
             size=44, bold=True, color=WHITE, font=F_HEAD, line_spacing=1.1)

    # акцентная линия
    add_rect(s, Inches(2.5), Inches(5.6), Inches(0.9), Inches(0.06), fill=NEON)
    add_text(s, Inches(3.5), Inches(5.45), Inches(7), Inches(0.4),
             "Михаил Илларионович Кутузов,  1812",
             size=14, color=WHITE, font=F_BODY)

    # подпись-росчерк
    add_text(s, Inches(2.5), Inches(6.0), Inches(6), Inches(0.6),
             "М. Кутузовъ", size=34, color=NEON, font=F_HAND, italic=True)

    set_transition(s, "fade", duration_ms=1000)


# =========================================================
#             СБОРКА
# =========================================================

slide_title()
slide_problem()
slide_break(3, "a", icon_map)
slide_break(4, "b", icon_clock)
slide_break(5, "c", icon_snow)
slide_solution()
slide_case()
slide_quote()
slide_legacy()
slide_final()

out = "Kutuzov.pptx"
prs.save(out)
print("OK:", out)
